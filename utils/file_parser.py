#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文件解析模块
支持Excel、CSV、Word文档的客户数据导入
"""

import os
import io
import pandas as pd
from utils.validators import sanitize_text

# 列名映射（支持多种命名方式）
COLUMN_MAPPINGS = {
    'customer_name': ['客户', '客户名称', '公司名称', 'company', 'customer', '客户名'],
    'country': ['国家', '国家/地区', 'country', '地区'],
    'address': ['地址', 'address', '公司地址'],
    'website': ['官网', '网站', 'website', '网址', '公司网站'],
    'company_info': ['公司信息和客户主要产品', '公司信息', '产品', 'company_info', '主营业务'],
    'supplier': ['供应商', 'supplier'],
    'supplier_info': ['供应商信息', 'supplier_info'],
    'customs_data': ['海关数据购买产品名', '海关数据', 'customs_data'],
    'logistics_info': ['物流信息', 'logistics_info'],
}

EMAIL_COLUMN_NAMES = ['客户邮箱', '邮箱', 'email', '邮件']
LINKEDIN_COLUMN_NAMES = ['领英邮箱', 'linkedin', '领英']


def detect_columns(df):
    """
    自动检测DataFrame中的列名映射
    返回: {standard_name: actual_column_name}
    """
    detected = {}
    df_columns = list(df.columns)
    df_columns_lower = [str(c).strip().lower() for c in df_columns]

    for standard_name, possible_names in COLUMN_MAPPINGS.items():
        for possible in possible_names:
            possible_lower = possible.lower()
            if possible_lower in df_columns_lower:
                idx = df_columns_lower.index(possible_lower)
                detected[standard_name] = df_columns[idx]
                break

    # 检测邮箱列
    for col in df_columns:
        col_lower = str(col).strip().lower()
        if col_lower in [e.lower() for e in EMAIL_COLUMN_NAMES]:
            detected['email_column'] = col
            break

    # 检测领英邮箱列
    for col in df_columns:
        col_lower = str(col).strip().lower()
        if col_lower in [e.lower() for e in LINKEDIN_COLUMN_NAMES]:
            detected['linkedin_column'] = col
            break

    return detected


def parse_emails(email_text):
    """
    从文本中解析多个邮箱，同时提取联系人姓名和职位
    支持多种格式：
      - 纯邮箱列表（逗号/分号/换行分隔）
      - 姓名：XXX\n职位：XXX\n邮箱：xxx@xxx.com
      - xxx@xxx.com=Name, Title
      - Name\nxxx@xxx.com
    返回: [(email_address, email_type, contact_name, job_title), ...]
    """
    if not email_text or pd.isna(email_text):
        return []

    import re

    email_text = str(email_text).strip()
    if not email_text or email_text.lower() in ('nan', 'none', 'null', '-'):
        return []

    results = []
    email_pattern = re.compile(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}')

    # 先尝试提取结构化数据（姓名+职位+邮箱）
    # 模式1: 姓名：XXX\n职位：XXX\n邮箱：xxx@xxx.com
    # 使用 [\n\r\u3000\s]+ 匹配换行和中文全角空格
    # 支持一个联系人对应多个邮箱，用中文顿号、逗号或英文逗号分隔
    pattern1 = r'姓名[：:]\s*([^\n]+)[\n\r\u3000\s]+职位[：:]\s*([^\n]*)[\n\r\u3000\s]+邮箱[：:]\s*([^\n]+)'
    for m in re.finditer(pattern1, email_text):
        name = m.group(1).strip()
        title = m.group(2).strip()
        email_field = m.group(3).strip()
        if name in ['-', '无姓名', '无', '']:
            name = None
        if title in ['-', '无', '']:
            title = None
        # 分割多个邮箱（支持中文顿号、英文逗号、分号、空格分隔）
        emails_in_field = re.split(r'[、,;\s]+', email_field)
        for email in emails_in_field:
            email = email.strip()
            if email_pattern.match(email):
                email_type = _detect_email_type(email)
                results.append((email, email_type, name, title))

    # 模式2: xxx@xxx.com=Name, Title
    pattern2 = r'([^\n\s=,]+@[^\n\s=,]+)=([^,\n]+)'
    for m in re.finditer(pattern2, email_text):
        email = m.group(1).strip()
        info = m.group(2).strip()
        name_title = re.match(r'^([^，,]+)[，,]\s*(.+)$', info)
        if name_title:
            name = name_title.group(1).strip()
            title = name_title.group(2).strip()
        else:
            name = info
            title = None
        email_type = _detect_email_type(email)
        results.append((email, email_type, name, title))

    # 模式3: Name\nxxx@xxx.com (多行格式) — 放宽正则，允许中间含额外单词
    lines = email_text.split('\n')
    for i in range(len(lines) - 1):
        line1 = lines[i].strip()
        line2 = lines[i + 1].strip()
        # 匹配 "First Last" 或 "First Middle Last" 或 "Tessa Bosche I Petromax"
        if re.match(r'^[A-Z][a-z]+(?:\s+[A-Za-z]+)+$', line1) and email_pattern.match(line2):
            email_type = _detect_email_type(line2)
            results.append((line2, email_type, line1, None))

    # 模式5: Name（职位）\nemail — 中文括号标注职位
    pattern5 = r'([A-Z][a-z]+(?:\s+[A-Za-z]+)*)[（(]([^)）]+)[)）][\n\r]+([a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})'
    for m in re.finditer(pattern5, email_text):
        name = m.group(1).strip()
        title = m.group(2).strip()
        email = m.group(3).strip()
        email_type = _detect_email_type(email)
        results.append((email, email_type, name, title))

    # 模式6: Name Title=email — 等号前是姓名+职位（反转格式）
    pattern6 = r'([A-Z][a-z]+(?:\s+[A-Za-z.]+)*)\s*=\s*([a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})'
    for m in re.finditer(pattern6, email_text):
        name_part = m.group(1).strip()
        email = m.group(2).strip()
        # 尝试从 name_part 分离姓名和职位
        nt = re.match(r'^([A-Z][a-z]+(?:\s+[A-Z][a-z]+)+)\s+([A-Z][a-z]+(?:\s+[A-Za-z]+)*)$', name_part)
        if nt:
            name = nt.group(1).strip()
            title = nt.group(2).strip()
        else:
            name = name_part
            title = None
        email_type = _detect_email_type(email)
        results.append((email, email_type, name, title))

    # 模式7: "Name" <email> 或 Name <email> — RFC 5322 标准格式
    pattern7 = r'"?([^"<>\n]+)"?\s*<\s*([a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})\s*>'
    for m in re.finditer(pattern7, email_text):
        name = m.group(1).strip()
        email = m.group(2).strip()
        # 清理姓名中的转义字符
        name = re.sub(r'\\(.)', r'\1', name)
        if name and not re.match(r'^[a-zA-Z0-9._%+-]+$', name):  # 排除纯邮箱前缀
            email_type = _detect_email_type(email)
            results.append((email, email_type, name, None))

    # 模式8: 缩写姓名 X. Surname 或 X Surname\nemail
    for i in range(len(lines) - 1):
        line1 = lines[i].strip()
        line2 = lines[i + 1].strip()
        if re.match(r'^[A-Z]\.?\s+[A-Za-z]+$', line1) and email_pattern.match(line2):
            email_type = _detect_email_type(line2)
            results.append((line2, email_type, line1, None))

    # 模式9: Name\nTitle\nemail 三行格式
    for i in range(len(lines) - 2):
        line1 = lines[i].strip()
        line2 = lines[i + 1].strip()
        line3 = lines[i + 2].strip()
        if re.match(r'^[A-Z][a-z]+(?:\s+[A-Za-z]+)*$', line1) and \
           email_pattern.match(line3) and \
           not email_pattern.match(line2) and \
           not re.match(r'^[A-Z][a-z]+(?:\s+[A-Za-z]+)*$', line2):
            email_type = _detect_email_type(line3)
            results.append((line3, email_type, line1, line2))

    # 模式10: 简单邮箱列表（无姓名信息）
    # 提取所有尚未被结构化模式捕获的邮箱
    found_emails = {r[0].lower() for r in results}
    for m in email_pattern.finditer(email_text):
        email = m.group(0)
        if email.lower() not in found_emails:
            email_type = _detect_email_type(email)
            results.append((email, email_type, None, None))

    return results


def _detect_email_type(email):
    """根据邮箱地址判断类型（使用新的邮箱分类器）"""
    from utils.email_classifier import classify_email
    email_type, _ = classify_email(email)
    return email_type


def parse_email_list(email_text):
    """
    从纯邮箱列表文本中解析邮箱（每行一个邮箱，无姓名信息）
    从邮箱前缀推断姓名，判断邮箱类型
    
    返回: [(email_address, email_type, contact_name, job_title), ...]
    """
    if not email_text or pd.isna(email_text):
        return []
    
    import re
    
    email_text = str(email_text).strip()
    if not email_text or email_text.lower() in ('nan', 'none', 'null', '-'):
        return []
    
    results = []
    email_pattern = re.compile(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}')
    
    # 按换行符、逗号、分号分隔
    lines = re.split(r'[\n\r,;]+', email_text)
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # 提取邮箱
        match = email_pattern.search(line)
        if match:
            email = match.group(0)
            email_type = _detect_email_type(email)
            
            # 从邮箱前缀推断姓名（使用新的邮箱分类器）
            from utils.email_classifier import infer_name_from_prefix
            prefix = email.split('@')[0]
            contact_name = infer_name_from_prefix(prefix)
            
            results.append((email, email_type, contact_name, None))
    
    return results


def _infer_name_from_prefix(prefix):
    """从邮箱前缀推断姓名"""
    import re
    
    # 清理前缀
    prefix = prefix.strip().lower()
    
    # 如果前缀是公共邮箱关键词，不推断姓名
    public_keywords = [
        'info', 'sales', 'support', 'contact', 'admin', 'hello', 'team',
        'service', 'help', 'marketing', 'office', 'general', 'enquiries',
        'inquiry', 'business', 'customerservice', 'feedback', 'hr', 'careers',
        'jobs', 'press', 'media', 'partners', 'abuse', 'webmaster', 'postmaster',
        'hostmaster', 'noc', 'security', 'billing', 'account', 'accounts',
        'finance', 'legal', 'privacy', 'recruitment', 'care', 'global-hr',
        'global.partnerships', 'customerservicetz', 'customerserviceug',
        'batterymasters', 'rvice', 'ervice'
    ]
    
    if prefix in public_keywords:
        return ''
    
    # 模式1: firstname.lastname (john.smith)
    if '.' in prefix:
        parts = prefix.split('.')
        if len(parts) == 2 and len(parts[0]) > 1 and len(parts[1]) > 1:
            first = parts[0].capitalize()
            last = parts[1].capitalize()
            return f"{first} {last}"
    
    # 模式2: firstname_lastname (john_smith)
    if '_' in prefix:
        parts = prefix.split('_')
        if len(parts) == 2 and len(parts[0]) > 1 and len(parts[1]) > 1:
            first = parts[0].capitalize()
            last = parts[1].capitalize()
            return f"{first} {last}"
    
    # 模式3: firstnamelastname (johnsmith) - 尝试分割
    if len(prefix) > 3 and prefix.isalpha():
        # 尝试找到名字和姓氏的分界（大写字母或常见名字模式）
        match = re.match(r'^([a-z]+)([a-z]+)$', prefix)
        if match:
            first, last = match.group(1), match.group(2)
            if len(first) >= 2 and len(last) >= 2:
                return f"{first.capitalize()} {last.capitalize()}"
    
    # 模式4: 首字母+姓氏 (jsmith, j.smith)
    match = re.match(r'^([a-z])\.?([a-z]+)$', prefix)
    if match:
        first_initial = match.group(1).upper()
        last = match.group(2).capitalize()
        return f"{first_initial}. {last}"
    
    # 无法推断，返回空
    return ''


def row_to_customer(row, column_map):
    """
    将DataFrame行转换为标准客户数据结构
    """
    customer = {
        'customer_name': '',
        'country': '',
        'address': '',
        'website': '',
        'company_info': '',
        'supplier': '',
        'supplier_info': '',
        'customs_data': '',
        'logistics_info': '',
        'emails': []
    }

    # 提取基本字段
    for field, col_name in column_map.items():
        if field in ['email_column', 'linkedin_column']:
            continue
        if col_name in row.index:
            value = row[col_name]
            if not pd.isna(value):
                customer[field] = sanitize_text(str(value))

    # 如果提供了 website 但缺少 http(s) 前缀，自动补全
    if customer['website'] and not customer['website'].startswith(('http://', 'https://')):
        customer['website'] = 'https://' + customer['website']

    # 提取邮箱（含联系人姓名和职位）
    all_emails = []

    # 客户邮箱
    if 'email_column' in column_map:
        col = column_map['email_column']
        if col in row.index:
            emails = parse_emails(row[col])
            for email, email_type, contact_name, job_title in emails:
                all_emails.append({
                    'email_address': email,
                    'email_type': email_type,
                    'contact_name': contact_name,
                    'job_title': job_title,
                    'source': 'customer_email'
                })

    # 领英邮箱
    if 'linkedin_column' in column_map:
        col = column_map['linkedin_column']
        if col in row.index:
            emails = parse_emails(row[col])
            for email, email_type, contact_name, job_title in emails:
                all_emails.append({
                    'email_address': email,
                    'email_type': 'linkedin',
                    'contact_name': contact_name,
                    'job_title': job_title,
                    'source': 'linkedin'
                })

    customer['emails'] = all_emails

    return customer


def parse_excel(file_path_or_obj):
    """
    解析Excel文件
    返回: ([customer_dict, ...], error_message or None)
    """
    try:
        df = pd.read_excel(file_path_or_obj, dtype=str)
        return _parse_dataframe(df)
    except Exception as e:
        return [], f"Excel解析失败: {str(e)}"


def parse_csv(file_path_or_obj):
    """
    解析CSV文件（自动检测编码）
    返回: ([customer_dict, ...], error_message or None)
    """
    encodings = ['utf-8', 'gbk', 'gb2312', 'utf-8-sig']

    for encoding in encodings:
        try:
            df = pd.read_csv(file_path_or_obj, dtype=str, encoding=encoding)
            return _parse_dataframe(df)
        except UnicodeDecodeError:
            if hasattr(file_path_or_obj, 'seek'):
                file_path_or_obj.seek(0)
            continue
        except Exception as e:
            return [], f"CSV解析失败: {str(e)}"

    return [], "CSV编码无法识别，请使用UTF-8或GBK编码"


def parse_word(file_path_or_obj):
    """
    解析Word文档（.docx）
    优先读取表格，无表格则尝试段落文本
    返回: ([customer_dict, ...], error_message or None)
    """
    try:
        from docx import Document
    except ImportError:
        return [], "未安装python-docx库，无法解析Word文档"

    try:
        if isinstance(file_path_or_obj, str):
            doc = Document(file_path_or_obj)
        else:
            doc = Document(file_path_or_obj)

        # 优先读取表格
        if doc.tables:
            all_customers = []
            for table in doc.tables:
                if len(table.rows) < 2:
                    continue

                # 提取表头
                headers = [cell.text.strip() for cell in table.rows[0].cells]

                # 提取数据行
                data = []
                for row in table.rows[1:]:
                    row_data = [cell.text.strip() for cell in row.cells]
                    if any(row_data):  # 跳过空行
                        data.append(row_data)

                if data:
                    df = pd.DataFrame(data, columns=headers)
                    customers, error = _parse_dataframe(df)
                    if error:
                        return [], error
                    all_customers.extend(customers)

            if all_customers:
                return all_customers, None

        # 无表格或表格为空，尝试段落文本（降级方案）
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        if paragraphs:
            # 尝试从段落中提取结构化数据
            return _parse_paragraphs(paragraphs)

        return [], "Word文档中未找到可导入的数据"

    except Exception as e:
        return [], f"Word解析失败: {str(e)}"


def _parse_dataframe(df):
    """
    解析DataFrame为客户数据列表
    """
    if df.empty:
        return [], "文件为空"

    # 检测列名
    column_map = detect_columns(df)

    if 'customer_name' not in column_map:
        return [], f"无法识别客户名称列。支持的列名: {', '.join(COLUMN_MAPPINGS['customer_name'])}"

    customers = []
    for _, row in df.iterrows():
        customer = row_to_customer(row, column_map)
        if customer['customer_name']:
            customers.append(customer)

    return customers, None


def _parse_paragraphs(paragraphs):
    """
    从段落文本中提取客户数据（降级方案）
    """
    import re

    customers = []
    current_customer = None

    for para in paragraphs:
        # 尝试匹配 "客户: XXX" 或 "公司名称: XXX" 格式
        name_match = re.search(r'(?:客户|公司名称|公司)[：:]\s*(.+)', para)
        if name_match:
            if current_customer:
                customers.append(current_customer)
            current_customer = {
                'customer_name': name_match.group(1).strip(),
                'country': '', 'address': '', 'website': '',
                'company_info': '', 'supplier': '',
                'supplier_info': '', 'customs_data': '',
                'logistics_info': '', 'emails': []
            }
            continue

        if current_customer:
            # 尝试匹配邮箱
            emails = parse_emails(para)
            for email, email_type, contact_name, job_title in emails:
                current_customer['emails'].append({
                    'email_address': email,
                    'email_type': email_type,
                    'contact_name': contact_name,
                    'job_title': job_title,
                    'source': 'customer_email'
                })

            # 尝试匹配网站
            website_match = re.search(r'(?:官网|网站|网址)[：:]\s*(https?://\S+)', para)
            if website_match:
                current_customer['website'] = website_match.group(1)

    if current_customer:
        customers.append(current_customer)

    return customers, None if customers else "无法从段落文本中提取客户数据"


def parse_file(file_path, filename=None):
    """
    根据文件扩展名自动选择解析器
    返回: ([customer_dict, ...], error_message or None)
    """
    if filename is None:
        filename = os.path.basename(file_path)

    ext = os.path.splitext(filename)[1].lower()

    if ext in ('.xlsx', '.xls'):
        return parse_excel(file_path)
    elif ext == '.csv':
        return parse_csv(file_path)
    elif ext == '.docx':
        return parse_word(file_path)
    elif ext == '.pdf':
        return parse_pdf(file_path)
    elif ext in ('.pptx', '.ppt'):
        return parse_pptx(file_path)
    else:
        return [], f"不支持的文件格式: {ext}"


def parse_pdf(file_path_or_obj):
    """
    解析 PDF 文件
    优先读取表格，否则提取全文交给 AI 分析
    返回: ([customer_dict, ...], error_message or None)
    """
    try:
        import pdfplumber
    except ImportError:
        return [], "未安装 pdfplumber 库，无法解析 PDF 文档"

    try:
        if isinstance(file_path_or_obj, str):
            pdf = pdfplumber.open(file_path_or_obj)
        else:
            # BytesIO 对象
            pdf = pdfplumber.open(file_path_or_obj)

        # 优先从表格提取结构化数据
        all_customers = []
        for page in pdf.pages:
            if page.extract_table():
                for table in page.extract_tables():
                    if len(table) < 2:
                        continue
                    headers = [str(cell).strip() if cell else '' for cell in table[0]]
                    data = []
                    for row in table[1:]:
                        row_data = [str(cell).strip() if cell else '' for cell in row]
                        if any(row_data):
                            data.append(row_data)
                    if data:
                        df = pd.DataFrame(data, columns=headers)
                        customers, error = _parse_dataframe(df)
                        if error:
                            continue
                        all_customers.extend(customers)

        if all_customers:
            pdf.close()
            return all_customers, None

        # 无表格：提取全文，返回纯文本用于 AI 分析
        all_text = []
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                all_text.append(text)
        pdf.close()

        full_text = '\n'.join(all_text).strip()
        if not full_text:
            return [], "PDF 文件为空或无法提取文字"

        # 尝试从文本中提取结构化数据
        paragraphs = [p.strip() for p in full_text.split('\n') if p.strip()]
        return _parse_paragraphs(paragraphs)

    except Exception as e:
        return [], f"PDF 解析失败: {str(e)}"


def parse_pptx(file_path_or_obj):
    """
    解析 PowerPoint 文件（.pptx）
    提取所有幻灯片文本，尝试结构化解析
    返回: ([customer_dict, ...], error_message or None)
    """
    try:
        from pptx import Presentation
    except ImportError:
        return [], "未安装 python-pptx 库，无法解析 PowerPoint 文档"

    try:
        if isinstance(file_path_or_obj, str):
            prs = Presentation(file_path_or_obj)
        else:
            prs = Presentation(file_path_or_obj)

        # 收集所有幻灯片文本
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            all_text.append(text)
                # 提取表格
                if shape.has_table:
                    table = shape.table
                    if len(table.rows) >= 2:
                        headers = [cell.text.strip() for cell in table.rows[0].cells]
                        data = []
                        for row in table.rows[1:]:
                            row_data = [cell.text.strip() for cell in row.cells]
                            if any(row_data):
                                data.append(row_data)
                        if data:
                            df = pd.DataFrame(data, columns=headers)
                            customers, _ = _parse_dataframe(df)
                            # 如果表格解析成功，直接返回
                            if customers:
                                return customers, None

        full_text = '\n'.join(all_text).strip()
        if not full_text:
            return [], "PowerPoint 文件中未找到文本内容"

        # 尝试从段落中提取
        paragraphs = [p.strip() for p in full_text.split('\n') if p.strip()]
        return _parse_paragraphs(paragraphs)

    except Exception as e:
        return [], f"PowerPoint 解析失败: {str(e)}"
