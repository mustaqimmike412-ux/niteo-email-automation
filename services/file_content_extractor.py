#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文件内容提取模块
从各种文件格式中提取纯文本内容，供 AI 分析使用
支持: Excel (.xlsx, .xls), CSV, Word (.docx), PDF, TXT, HTML (.html, .htm), PPT (.pptx)
"""

import os
import io
import pandas as pd
from typing import Tuple, Optional


def extract_from_file(file_path: str, filename: str = None) -> Tuple[str, str]:
    """从文件中提取纯文本内容（兼容旧接口）"""
    return extract_text_from_file(file_path, filename)


def extract_text_from_file(file_path: str, filename: str = None) -> Tuple[str, str]:
    """
    从文件中提取纯文本内容

    Args:
        file_path: 文件路径
        filename: 文件名（用于判断文件类型）

    Returns:
        (text_content, error_message)
        text_content: 提取的文本内容
        error_message: 错误信息（成功为 None）
    """
    if filename is None:
        filename = os.path.basename(file_path)

    ext = os.path.splitext(filename)[1].lower()

    try:
        if ext in ('.xlsx', '.xls'):
            return extract_from_excel(file_path)
        elif ext == '.csv':
            return extract_from_csv(file_path)
        elif ext == '.docx':
            return extract_from_word(file_path)
        elif ext == '.pdf':
            return extract_from_pdf(file_path)
        elif ext == '.txt':
            return extract_from_txt(file_path)
        elif ext in ('.html', '.htm'):
            return extract_from_html(file_path)
        elif ext == '.pptx':
            return extract_from_ppt(file_path)
        else:
            return "", f"不支持的文件格式: {ext}"
    except Exception as e:
        return "", f"文件提取失败: {str(e)}"


def extract_from_excel(file_path: str) -> Tuple[str, str]:
    """从 Excel 文件中提取文本内容"""
    try:
        # 读取所有 sheet
        xl = pd.ExcelFile(file_path)
        all_text = []

        for sheet_name in xl.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name, dtype=str)
            if df.empty:
                continue

            all_text.append(f"=== Sheet: {sheet_name} ===")

            # 添加列名
            headers = [str(col) for col in df.columns]
            all_text.append(" | ".join(headers))
            all_text.append("-" * 50)

            # 添加数据行
            for _, row in df.iterrows():
                row_values = []
                for val in row.values:
                    if pd.isna(val):
                        row_values.append("")
                    else:
                        row_values.append(str(val))
                # 只添加非空行
                if any(v.strip() for v in row_values):
                    all_text.append(" | ".join(row_values))

            all_text.append("")  # 空行分隔

        return "\n".join(all_text), None
    except Exception as e:
        return "", f"Excel 提取失败: {str(e)}"


def extract_from_csv(file_path: str) -> Tuple[str, str]:
    """从 CSV 文件中提取文本内容"""
    try:
        encodings = ['utf-8', 'gbk', 'gb2312', 'utf-8-sig']

        for encoding in encodings:
            try:
                df = pd.read_csv(file_path, dtype=str, encoding=encoding)
                if df.empty:
                    return "", "CSV 文件为空"

                all_text = []
                headers = [str(col) for col in df.columns]
                all_text.append(" | ".join(headers))
                all_text.append("-" * 50)

                for _, row in df.iterrows():
                    row_values = []
                    for val in row.values:
                        if pd.isna(val):
                            row_values.append("")
                        else:
                            row_values.append(str(val))
                    if any(v.strip() for v in row_values):
                        all_text.append(" | ".join(row_values))

                return "\n".join(all_text), None
            except UnicodeDecodeError:
                continue

        return "", "CSV 编码无法识别"
    except Exception as e:
        return "", f"CSV 提取失败: {str(e)}"


def extract_from_word(file_path: str) -> Tuple[str, str]:
    """从 Word 文件中提取文本内容"""
    try:
        from docx import Document
    except ImportError:
        return "", "未安装 python-docx 库"

    try:
        doc = Document(file_path)
        all_text = []

        # 提取段落
        for para in doc.paragraphs:
            if para.text.strip():
                all_text.append(para.text.strip())

        # 提取表格
        if doc.tables:
            all_text.append("\n=== 表格内容 ===")
            for table_idx, table in enumerate(doc.tables):
                all_text.append(f"\n--- 表格 {table_idx + 1} ---")
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    if any(row_text):
                        all_text.append(" | ".join(row_text))

        return "\n".join(all_text), None
    except Exception as e:
        return "", f"Word 提取失败: {str(e)}"


def extract_from_pdf(file_path: str) -> Tuple[str, str]:
    """从 PDF 文件中提取文本内容"""
    try:
        import PyPDF2
    except ImportError:
        return "", "未安装 PyPDF2 库"

    try:
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            all_text = []

            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text and text.strip():
                    all_text.append(f"=== Page {page_num + 1} ===")
                    all_text.append(text.strip())

        return "\n\n".join(all_text), None
    except Exception as e:
        return "", f"PDF 提取失败: {str(e)}"


def extract_from_txt(file_path: str) -> Tuple[str, str]:
    """从 TXT 文件中提取文本内容"""
    try:
        encodings = ['utf-8', 'gbk', 'gb2312']

        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                return content, None
            except UnicodeDecodeError:
                continue

        return "", "TXT 编码无法识别"
    except Exception as e:
        return "", f"TXT 读取失败: {str(e)}"


def extract_from_bytes(file_bytes: bytes, filename: str) -> Tuple[str, str]:
    """
    从字节数据中提取文本内容

    Args:
        file_bytes: 文件字节数据
        filename: 文件名

    Returns:
        (text_content, error_message)
    """
    ext = os.path.splitext(filename)[1].lower()

    try:
        if ext in ('.xlsx', '.xls'):
            return _extract_excel_from_bytes(file_bytes)
        elif ext == '.csv':
            return _extract_csv_from_bytes(file_bytes)
        elif ext == '.docx':
            return _extract_word_from_bytes(file_bytes)
        elif ext == '.pdf':
            return _extract_pdf_from_bytes(file_bytes)
        elif ext == '.txt':
            return _extract_txt_from_bytes(file_bytes)
        elif ext in ('.html', '.htm'):
            return _extract_html_from_bytes(file_bytes)
        elif ext == '.pptx':
            return _extract_ppt_from_bytes(file_bytes)
        else:
            return "", f"不支持的文件格式: {ext}"
    except Exception as e:
        return "", f"文件提取失败: {str(e)}"


def _extract_excel_from_bytes(file_bytes: bytes) -> Tuple[str, str]:
    """从 Excel 字节数据中提取文本"""
    try:
        xl = pd.ExcelFile(io.BytesIO(file_bytes))
        all_text = []

        for sheet_name in xl.sheet_names:
            df = pd.read_excel(io.BytesIO(file_bytes), sheet_name=sheet_name, dtype=str)
            if df.empty:
                continue

            all_text.append(f"=== Sheet: {sheet_name} ===")
            headers = [str(col) for col in df.columns]
            all_text.append(" | ".join(headers))
            all_text.append("-" * 50)

            for _, row in df.iterrows():
                row_values = []
                for val in row.values:
                    if pd.isna(val):
                        row_values.append("")
                    else:
                        row_values.append(str(val))
                if any(v.strip() for v in row_values):
                    all_text.append(" | ".join(row_values))

            all_text.append("")

        return "\n".join(all_text), None
    except Exception as e:
        return "", f"Excel 提取失败: {str(e)}"


def _extract_csv_from_bytes(file_bytes: bytes) -> Tuple[str, str]:
    """从 CSV 字节数据中提取文本"""
    encodings = ['utf-8', 'gbk', 'gb2312', 'utf-8-sig']

    for encoding in encodings:
        try:
            df = pd.read_csv(io.BytesIO(file_bytes), dtype=str, encoding=encoding)
            if df.empty:
                return "", "CSV 文件为空"

            all_text = []
            headers = [str(col) for col in df.columns]
            all_text.append(" | ".join(headers))
            all_text.append("-" * 50)

            for _, row in df.iterrows():
                row_values = []
                for val in row.values:
                    if pd.isna(val):
                        row_values.append("")
                    else:
                        row_values.append(str(val))
                if any(v.strip() for v in row_values):
                    all_text.append(" | ".join(row_values))

            return "\n".join(all_text), None
        except UnicodeDecodeError:
            continue

    return "", "CSV 编码无法识别"


def _extract_word_from_bytes(file_bytes: bytes) -> Tuple[str, str]:
    """从 Word 字节数据中提取文本"""
    try:
        from docx import Document
    except ImportError:
        return "", "未安装 python-docx 库"

    try:
        doc = Document(io.BytesIO(file_bytes))
        all_text = []

        for para in doc.paragraphs:
            if para.text.strip():
                all_text.append(para.text.strip())

        if doc.tables:
            all_text.append("\n=== 表格内容 ===")
            for table_idx, table in enumerate(doc.tables):
                all_text.append(f"\n--- 表格 {table_idx + 1} ---")
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    if any(row_text):
                        all_text.append(" | ".join(row_text))

        return "\n".join(all_text), None
    except Exception as e:
        return "", f"Word 提取失败: {str(e)}"


def _extract_pdf_from_bytes(file_bytes: bytes) -> Tuple[str, str]:
    """从 PDF 字节数据中提取文本"""
    try:
        import PyPDF2
    except ImportError:
        return "", "未安装 PyPDF2 库"

    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        all_text = []

        for page_num, page in enumerate(pdf_reader.pages):
            text = page.extract_text()
            if text and text.strip():
                all_text.append(f"=== Page {page_num + 1} ===")
                all_text.append(text.strip())

        return "\n\n".join(all_text), None
    except Exception as e:
        return "", f"PDF 提取失败: {str(e)}"


def _extract_txt_from_bytes(file_bytes: bytes) -> Tuple[str, str]:
    """从 TXT 字节数据中提取文本"""
    encodings = ['utf-8', 'gbk', 'gb2312']

    for encoding in encodings:
        try:
            content = file_bytes.decode(encoding)
            return content, None
        except UnicodeDecodeError:
            continue

    return "", "TXT 编码无法识别"


def extract_from_html(file_path: str) -> Tuple[str, str]:
    """从 HTML 文件中提取纯文本内容"""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return "", "未安装 beautifulsoup4 库"

    try:
        encodings = ['utf-8', 'gbk', 'gb2312']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    html_content = f.read()
                break
            except UnicodeDecodeError:
                continue
        else:
            return "", "HTML 编码无法识别"

        soup = BeautifulSoup(html_content, 'html.parser')

        # 移除 script 和 style 标签
        for tag in soup(['script', 'style']):
            tag.decompose()

        text = soup.get_text(separator='\n', strip=True)

        return text, None
    except Exception as e:
        return "", f"HTML 提取失败: {str(e)}"


def _extract_html_from_bytes(file_bytes: bytes) -> Tuple[str, str]:
    """从 HTML 字节数据中提取纯文本"""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return "", "未安装 beautifulsoup4 库"

    try:
        encodings = ['utf-8', 'gbk', 'gb2312']
        for encoding in encodings:
            try:
                html_content = file_bytes.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            return "", "HTML 编码无法识别"

        soup = BeautifulSoup(html_content, 'html.parser')

        # 移除 script 和 style 标签
        for tag in soup(['script', 'style']):
            tag.decompose()

        text = soup.get_text(separator='\n', strip=True)

        return text, None
    except Exception as e:
        return "", f"HTML 提取失败: {str(e)}"


def extract_from_ppt(file_path: str) -> Tuple[str, str]:
    """从 PPT 文件中提取文本内容"""
    try:
        from pptx import Presentation
    except ImportError:
        return "", "未安装 python-pptx 库"

    try:
        prs = Presentation(file_path)
        all_text = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_texts = []
            slide_texts.append(f"=== 幻灯片 {slide_idx + 1} ===")

            # 提取 shapes 中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

                # 提取表格中的文本
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        if any(row_text):
                            slide_texts.append(" | ".join(row_text))

            if len(slide_texts) > 1:
                all_text.append("\n".join(slide_texts))

        return "\n\n".join(all_text), None
    except Exception as e:
        return "", f"PPT 提取失败: {str(e)}"


def _extract_ppt_from_bytes(file_bytes: bytes) -> Tuple[str, str]:
    """从 PPT 字节数据中提取文本"""
    try:
        from pptx import Presentation
    except ImportError:
        return "", "未安装 python-pptx 库"

    try:
        prs = Presentation(io.BytesIO(file_bytes))
        all_text = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_texts = []
            slide_texts.append(f"=== 幻灯片 {slide_idx + 1} ===")

            # 提取 shapes 中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

                # 提取表格中的文本
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        if any(row_text):
                            slide_texts.append(" | ".join(row_text))

            if len(slide_texts) > 1:
                all_text.append("\n".join(slide_texts))

        return "\n\n".join(all_text), None
    except Exception as e:
        return "", f"PPT 提取失败: {str(e)}"
